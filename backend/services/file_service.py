import os
import hashlib
import json
import time
from typing import List, Dict, Any, Tuple
from fastapi import UploadFile
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
try:
    import pytesseract
    from PIL import Image
    HAS_TESSERACT = True
except ImportError:
    HAS_TESSERACT = False

UPLOAD_DIR = "uploads"
OUTPUT_DIR = "outputs"

os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)

class FileService:
    @staticmethod
    async def save_files_and_hash(files: List[UploadFile]) -> Tuple[str, List[str]]:
        """Saves uploaded files and returns a combined SHA-256 hash and the file paths."""
        file_paths = []
        hasher = hashlib.sha256()
        
        # Sort files by name to ensure consistent hashing regardless of upload order
        sorted_files = sorted(files, key=lambda f: f.filename)
        
        for file in sorted_files:
            content = await file.read()
            hasher.update(content)
            
            # Save file
            file_path = os.path.join(UPLOAD_DIR, file.filename)
            with open(file_path, "wb") as f:
                f.write(content)
            file_paths.append(file_path)
            
            # Reset file cursor for future reads if needed
            await file.seek(0)
            
        combined_hash = hasher.hexdigest()
        return combined_hash, file_paths

    @staticmethod
    def get_cached_quiz(file_hash: str) -> Dict[str, Any]:
        """Checks if a generated quiz exists for the given hash within the last 24 hours."""
        cache_file = os.path.join(OUTPUT_DIR, f"{file_hash}.json")
        
        if os.path.exists(cache_file):
            file_age = time.time() - os.path.getmtime(cache_file)
            # 24 hours = 86400 seconds
            if file_age < 86400:
                with open(cache_file, "r", encoding="utf-8") as f:
                    return json.load(f)
        return None

    @staticmethod
    def save_to_cache(file_hash: str, quiz_data: Dict[str, Any]):
        """Saves generated quiz JSON to cache."""
        cache_file = os.path.join(OUTPUT_DIR, f"{file_hash}.json")
        with open(cache_file, "w", encoding="utf-8") as f:
            json.dump(quiz_data, f, ensure_ascii=False, indent=2)

    @staticmethod
    def extract_text_from_file(file_path: str, use_vision: bool) -> str:
        """Extracts text from a single file based on its extension."""
        ext = os.path.splitext(file_path)[1].lower()
        extracted_text = ""
        
        if ext == ".pdf":
            extracted_text = FileService._extract_pdf(file_path, use_vision)
        elif ext == ".docx":
            extracted_text = FileService._extract_docx(file_path)
        elif ext == ".pptx":
            extracted_text = FileService._extract_pptx(file_path)
        else:
            # Fallback for txt or other formats
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    extracted_text = f.read()
            except Exception:
                extracted_text = f"Unsupported or unreadable file format: {ext}"
                
        return extracted_text

    @staticmethod
    def _extract_pdf(file_path: str, use_vision: bool) -> str:
        text = ""
        try:
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_text = page.get_text("text")
                text += page_text + "\n"
                
                # If text is extremely short, it might be a scanned image
                if len(page_text.strip()) < 50 and not use_vision and HAS_TESSERACT:
                    pix = page.get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    ocr_text = pytesseract.image_to_string(img, lang="eng+hin+ara")
                    text += ocr_text + "\n"
                    
                # Note: If use_vision is True, we ideally want to send the images to the AI.
                # Since OpenRouter supports passing text blocks or image URLs, for Phase 1 
                # we'll extract what we can. A full Vision implementation would convert 
                # pages to base64 and send them in the AI prompt.
                
        except Exception as e:
            text = f"Error reading PDF: {str(e)}"
        return text

    @staticmethod
    def _extract_docx(file_path: str) -> str:
        text = ""
        try:
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        except Exception as e:
            text = f"Error reading DOCX: {str(e)}"
        return text

    @staticmethod
    def _extract_pptx(file_path: str) -> str:
        text = ""
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            text = f"Error reading PPTX: {str(e)}"
        return text
